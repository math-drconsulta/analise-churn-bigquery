"""
Gera apresentacao PowerPoint com os insights de churn.
Uso: uv add python-pptx && uv run python scripts/gerar_ppt.py
Saida: churn_insights.pptx
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Cores ──────────────────────────────────────────────────────────
AZUL_ESCURO = RGBColor(0x1B, 0x2A, 0x4A)
VERMELHO = RGBColor(0xC0, 0x39, 0x2B)
VERDE = RGBColor(0x27, 0xAE, 0x60)
LARANJA = RGBColor(0xE6, 0x7E, 0x22)
CINZA = RGBColor(0x7F, 0x8C, 0x8D)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xEC, 0xF0, 0xF1)
AZUL = RGBColor(0x29, 0x80, 0xB9)

# ── Ticket medio (ajustar conforme necessario) ────────────────────
TICKET = {"6": 600.0, "12": 1200.0}


# ── Helpers ────────────────────────────────────────────────────────
def add_slide(prs, layout_idx=6):
    """Adiciona slide em branco."""
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def add_title_bar(slide, text, subtitle=None):
    """Barra de titulo no topo do slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = AZUL_ESCURO
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.color.rgb = BRANCO
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    if subtitle:
        shape_sub = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.15), Inches(12), Inches(0.4)
        )
        tf_sub = shape_sub.text_frame
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = CINZA


def add_text(slide, left, top, width, height, text, size=14, bold=False, color=AZUL_ESCURO, align=PP_ALIGN.LEFT):
    """Adiciona caixa de texto."""
    txbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, size=13):
    """Lista com marcadores."""
    txbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = AZUL_ESCURO
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_kpi_box(slide, left, top, label, value, color=AZUL_ESCURO, delta=None, delta_color=VERDE):
    """Card de KPI."""
    width, height = 2.8, 1.3

    # Fundo
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CINZA_CLARO
    shape.line.fill.background()

    # Valor
    add_text(slide, left + 0.15, top + 0.1, width - 0.3, 0.6,
             value, size=24, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Label
    add_text(slide, left + 0.15, top + 0.7, width - 0.3, 0.3,
             label, size=11, color=CINZA, align=PP_ALIGN.CENTER)

    # Delta
    if delta:
        add_text(slide, left + 0.15, top + 1.0, width - 0.3, 0.25,
                 delta, size=10, color=delta_color, align=PP_ALIGN.CENTER)


def add_table(slide, left, top, width, data, col_widths=None):
    """Tabela simples."""
    rows = len(data)
    cols = len(data[0]) if rows > 0 else 0
    if rows == 0 or cols == 0:
        return

    table_shape = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top),
        Inches(width), Inches(rows * 0.4)
    )
    table = table_shape.table

    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = AZUL_ESCURO if i > 0 else BRANCO
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT

            if i == 0:  # Header
                cell.fill.solid()
                cell.fill.fore_color.rgb = AZUL_ESCURO
                p.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BRANCO if i % 2 == 1 else CINZA_CLARO

    if col_widths:
        for j, w in enumerate(col_widths):
            table.columns[j].width = Inches(w)

    return table


# ── Carregar dados ─────────────────────────────────────────────────
print("Carregando dados...")
df_uni = pd.read_csv("results/univariada.csv")
df_int = pd.read_csv("results/interacao_contrato_dep_cronico.csv")
df_fin = pd.read_csv("results/impacto_financeiro.csv")
df_fin["duracao"] = df_fin["duracao"].astype(str)
df_clv = pd.read_csv("results/clv_por_perfil.csv")
df_clv["duracao"] = df_clv["duracao"].astype(str)
df_sil = pd.read_csv("results/churn_silencioso_vs_ativo.csv")
df_vel = pd.read_csv("results/velocidade_churn.csv")
df_vel["duracao"] = df_vel["duracao"].astype(str)
df_wb = pd.read_csv("results/winback_reativacoes.csv")

# Metricas globais
total = int(df_uni[df_uni["dimensao"] == "plan_months_duration"]["total_contratos"].sum())
churners_n = int(df_uni[df_uni["dimensao"] == "plan_months_duration"]["churners"].sum())
retidos = total - churners_n
taxa = round(100 * churners_n / total, 1)

df_fin["ticket"] = df_fin["duracao"].map(TICKET)
receita_perdida = (df_fin["churners"] * df_fin["ticket"]).sum()
receita_total = (df_fin["total_contratos"] * df_fin["ticket"]).sum()

df_clv["ticket"] = df_clv["duracao"].map(TICKET)
df_clv["ticket_mensal"] = df_clv["ticket"] / df_clv["duracao"].astype(float)
df_clv["clv"] = df_clv["ticket_mensal"] * df_clv["meses_vida_estimados"]

# Drivers
df_ciclo = df_uni[df_uni["dimensao"] == "account_contract_number"]
cr_1o = float(df_ciclo[df_ciclo["segmento"] == "1o contrato"]["churn_rate"].values[0])
cr_2o = float(df_ciclo[df_ciclo["segmento"] == "2o+ contrato"]["churn_rate"].values[0])

df_dep = df_uni[df_uni["dimensao"] == "dependentes_faixa"]
cr_sem_dep = float(df_dep[df_dep["segmento"] == "0 (sem dep.)"]["churn_rate"].values[0])
cr_com_dep = float(df_dep[df_dep["segmento"] == "3+ dep."]["churn_rate"].values[0])

df_cron = df_uni[df_uni["dimensao"] == "titular_main_cronico_sn"]
cr_nao_cron = float(df_cron[df_cron["segmento"] == "N"]["churn_rate"].values[0])
cr_sim_cron = float(df_cron[df_cron["segmento"] == "S"]["churn_rate"].values[0])

# Interacao
df_int_sorted = df_int.sort_values("churn_rate", ascending=False)
pior = df_int_sorted.iloc[0]
melhor = df_int_sorted.iloc[-1]
gap = round(pior["churn_rate"] - melhor["churn_rate"], 1)

# Silencioso
total_sil = df_sil.groupby("tipo_desfecho")["total_contratos"].sum()
n_silencioso = int(total_sil.get("churn_silencioso", 0))
n_ativo = int(total_sil.get("churn_ativo", 0))
pct_sil = round(100 * n_silencioso / (n_silencioso + n_ativo), 0)

# CLV
clv_1o = df_clv[df_clv["ciclo"] == "1o"]["clv"].mean()
clv_2o = df_clv[df_clv["ciclo"] == "2o+"]["clv"].mean()

# Janela
df_vel_ativo = df_vel[df_vel["tipo_churn"] == "churn_ativo"]
janela_agg = df_vel_ativo.groupby("janela_saida")["total"].sum().reset_index()
janela_agg["pct"] = round(100 * janela_agg["total"] / janela_agg["total"].sum(), 1)
antecipado_pct = janela_agg[janela_agg["janela_saida"].isin([
    "A_90+_dias_antes", "B_31-90_dias_antes"
])]["pct"].sum()


# ══════════════════════════════════════════════════════════════════
# MONTAR PPT
# ══════════════════════════════════════════════════════════════════
print("Gerando apresentacao...")
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


# ── SLIDE 1: CAPA ──────────────────────────────────────────────────
slide = add_slide(prs)
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = AZUL_ESCURO
shape.line.fill.background()

add_text(slide, 1, 2.0, 11, 1.2,
         "Analise de Churn", size=44, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.2, 11, 0.8,
         "O que os numeros mostram", size=28, color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)
add_text(slide, 1, 4.5, 11, 0.5,
         "Contratos cartao de credito  |  Planos 6 e 12 meses  |  Ultimos 12 meses",
         size=16, color=CINZA, align=PP_ALIGN.CENTER)
add_text(slide, 1, 6.2, 11, 0.4,
         "dr.consulta  |  Dados & Analytics", size=14, color=CINZA, align=PP_ALIGN.CENTER)


# ── SLIDE 2: O TAMANHO DO PROBLEMA ────────────────────────────────
slide = add_slide(prs)
add_title_bar(slide, "1. O tamanho do problema")

add_kpi_box(slide, 0.5, 1.5, "Contratos analisados", f"{total:,}")
add_kpi_box(slide, 3.5, 1.5, "Nao renovaram", f"{churners_n:,}",
            color=VERMELHO, delta=f"{taxa}% de churn", delta_color=VERMELHO)
add_kpi_box(slide, 6.5, 1.5, "Receita nao renovada",
            f"R$ {receita_perdida/1e6:.1f}M", color=VERMELHO,
            delta=f"{100*receita_perdida/receita_total:.0f}% da receita", delta_color=VERMELHO)
add_kpi_box(slide, 9.5, 1.5, "Renovaram", f"{retidos:,}",
            color=VERDE, delta=f"{100-taxa:.1f}%", delta_color=VERDE)

add_text(slide, 0.5, 3.3, 12, 1.5,
         f"De cada 100 contratos que vencem, {taxa:.0f} nao renovam.\n\n"
         f"Isso representa R$ {receita_perdida:,.0f} em contratos que nao geraram "
         f"um proximo ciclo — {100*receita_perdida/receita_total:.0f}% de toda a receita do periodo.\n\n"
         f"Mas esse numero nao e homogeneo. Quando separamos por perfil, "
         f"o churn vai de {melhor['churn_rate']}% ate {pior['churn_rate']}%.",
         size=16)


# ── SLIDE 3: ONDE ESTA CONCENTRADO ────────────────────────────────
slide = add_slide(prs)
add_title_bar(slide, "2. O churn nao e uniforme",
              "Com apenas 3 variaveis, separamos perfis com gap de " + f"{gap} p.p.")

# Tabela de drivers
add_table(slide, 0.5, 1.7, 7, [
    ["Variavel", "Grupo de risco", "Churn", "Grupo protegido", "Churn", "Gap"],
    ["Ciclo do contrato", "1o contrato", f"{cr_1o}%", "2o+ contrato", f"{cr_2o}%",
     f"{cr_1o - cr_2o:+.1f} p.p."],
    ["Dependentes", "Sem dependentes", f"{cr_sem_dep}%", "3+ dependentes", f"{cr_com_dep}%",
     f"{cr_sem_dep - cr_com_dep:+.1f} p.p."],
    ["Condicao cronica", "Nao cronico", f"{cr_nao_cron}%", "Cronico", f"{cr_sim_cron}%",
     f"{cr_nao_cron - cr_sim_cron:+.1f} p.p."],
], col_widths=[2.0, 1.2, 0.7, 1.3, 0.7, 1.1])

# Extremos
add_text(slide, 8.0, 1.7, 4.8, 0.4,
         "Perfis extremos:", size=16, bold=True)

# Pior
shape_p = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(2.3), Inches(4.8), Inches(1.6)
)
shape_p.fill.solid()
shape_p.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
shape_p.line.fill.background()
add_text(slide, 8.2, 2.4, 4.4, 0.3,
         f"MAIOR RISCO — Churn {pior['churn_rate']}%", size=14, bold=True, color=VERMELHO)
add_text(slide, 8.2, 2.8, 4.4, 0.9,
         f"{pior['ciclo']}o contrato | {pior['dependentes']} | "
         f"{'Cronico' if pior['cronico']=='S' else 'Nao cronico'}\n"
         f"{int(pior['total_contratos']):,} contratos",
         size=12, color=AZUL_ESCURO)

# Melhor
shape_m = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(4.2), Inches(4.8), Inches(1.6)
)
shape_m.fill.solid()
shape_m.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
shape_m.line.fill.background()
add_text(slide, 8.2, 4.3, 4.4, 0.3,
         f"MENOR RISCO — Churn {melhor['churn_rate']}%", size=14, bold=True, color=VERDE)
add_text(slide, 8.2, 4.7, 4.4, 0.9,
         f"{melhor['ciclo']}o contrato | {melhor['dependentes']} | "
         f"{'Cronico' if melhor['cronico']=='S' else 'Nao cronico'}\n"
         f"{int(melhor['total_contratos']):,} contratos",
         size=12, color=AZUL_ESCURO)

add_text(slide, 0.5, 5.5, 12, 1.0,
         f"A diferenca entre esses perfis e de {gap} p.p. — usando apenas ciclo, "
         f"dependentes e condicao cronica.\n"
         f"Onde tem estrutura, tem alavanca.",
         size=15, color=AZUL_ESCURO)


# ── SLIDE 4: O QUE PROTEGE ───────────────────────────────────────
slide = add_slide(prs)
add_title_bar(slide, "3. O que os clientes que ficam tem em comum")

# 3 colunas
for i, (titulo, v_risco, v_prot, l_risco, l_prot, texto) in enumerate([
    ("Vinculo familiar",
     f"{cr_sem_dep}%", f"{cr_com_dep}%", "Sem dependentes", "3+ dependentes",
     "Quando o plano cuida de mais\npessoas, o custo de sair sobe.\nCada dep. reduz ~3-5 p.p."),
    ("Necessidade medica",
     f"{cr_nao_cron}%", f"{cr_sim_cron}%", "Nao cronico", "Cronico",
     "Pacientes com condicao cronica\ndependem do acompanhamento.\nNao e fidelidade — e necessidade."),
    ("Experiencia previa",
     f"{cr_1o}%", f"{cr_2o}%", "1o contrato", "2o+ contrato",
     "Quem ja renovou uma vez\ntem probabilidade muito maior\nde renovar de novo."),
]):
    left = 0.5 + i * 4.2
    add_text(slide, left, 1.5, 3.8, 0.4, titulo, size=18, bold=True, color=AZUL)

    add_text(slide, left, 2.1, 1.8, 0.3, l_risco, size=11, color=CINZA)
    add_text(slide, left, 2.4, 1.8, 0.5, v_risco, size=28, bold=True, color=VERMELHO)

    add_text(slide, left + 2.0, 2.1, 1.8, 0.3, l_prot, size=11, color=CINZA)
    add_text(slide, left + 2.0, 2.4, 1.8, 0.5, v_prot, size=28, bold=True, color=VERDE)

    add_text(slide, left, 3.2, 3.8, 1.5, texto, size=13, color=AZUL_ESCURO)


# ── SLIDE 5: ANATOMIA DA SAIDA ───────────────────────────────────
slide = add_slide(prs)
add_title_bar(slide, "4. Como os clientes saem",
              f"{pct_sil:.0f}% do churn e silencioso — o paciente nao pediu para sair")

add_kpi_box(slide, 0.5, 1.5, "Churn silencioso", f"{n_silencioso:,}",
            color=LARANJA, delta=f"{pct_sil:.0f}% do total", delta_color=LARANJA)
add_kpi_box(slide, 3.5, 1.5, "Churn ativo", f"{n_ativo:,}",
            color=VERMELHO, delta=f"{100-pct_sil:.0f}% do total", delta_color=VERMELHO)
add_kpi_box(slide, 6.5, 1.5, "Cancelam com 30+ dias", f"{antecipado_pct:.0f}%",
            color=AZUL, delta="dos ativos", delta_color=CINZA)

add_text(slide, 0.5, 3.3, 12, 1.0,
         f"A maioria dos churners nao decidiu sair. O contrato venceu sem renovacao:\n"
         f"cartao recusado, cobranca nao processada, ou nenhuma tentativa registrada.\n\n"
         f"Dos que pediram cancelamento, {antecipado_pct:.0f}% o fizeram com mais de 30 dias "
         f"de antecedencia — ha uma janela para entender o motivo.",
         size=15)

add_text(slide, 0.5, 5.0, 12, 1.0,
         "Perguntas que isso levanta:\n"
         f"  - Quantos dos {n_silencioso:,} silenciosos sao recuperaveis com "
         f"uma acao simples (lembrete, atualizacao de cartao)?\n"
         f"  - O que leva alguem a cancelar 90+ dias antes do vencimento?",
         size=14, color=CINZA)


# ── SLIDE 6: VALOR EM JOGO (CLV) ─────────────────────────────────
slide = add_slide(prs)
add_title_bar(slide, "5. Quanto cada perfil vale",
              "Nem todo churn custa igual")

add_kpi_box(slide, 0.5, 1.5, "CLV medio 1o contrato", f"R$ {clv_1o:,.0f}", color=VERMELHO)
add_kpi_box(slide, 3.5, 1.5, "CLV medio 2o+ contrato", f"R$ {clv_2o:,.0f}", color=VERDE)
add_kpi_box(slide, 6.5, 1.5, "Premio de sobrevivencia",
            f"R$ {clv_2o - clv_1o:,.0f}", color=AZUL,
            delta=f"{clv_2o/clv_1o:.1f}x mais valor", delta_color=AZUL)

# Top 5 CLV
top5 = df_clv.sort_values("clv", ascending=False).head(5)
table_data = [["Perfil", "Duracao", "Churn", "Vida (meses)", "CLV"]]
for _, r in top5.iterrows():
    perfil = f"{r['ciclo']} | {r['perfil_idade']} | {'Cron.' if r['cronico']=='S' else 'N.cron.'} | {r['tem_dependente']}"
    table_data.append([
        perfil, f"{r['duracao']}m", f"{r['churn_rate_pct']}%",
        f"{r['meses_vida_estimados']:.0f}", f"R$ {r['clv']:,.0f}"
    ])
add_text(slide, 0.5, 3.1, 6, 0.3, "Top 5 perfis por CLV:", size=14, bold=True)
add_table(slide, 0.5, 3.5, 8, table_data, col_widths=[3.0, 0.8, 0.8, 1.2, 1.5])

add_text(slide, 0.5, 5.8, 12, 0.8,
         f"Cada cliente que sobrevive ao 1o ciclo vale {clv_2o/clv_1o:.1f}x mais.\n"
         f"O primeiro ciclo e o filtro mais caro.",
         size=15)


# ── SLIDE 7: E SE? ───────────────────────────────────────────────
slide = add_slide(prs)
add_title_bar(slide, "6. E se o churn caisse 5 p.p.?")

reducao = 5
evitaveis = int(total * reducao / 100)
ticket_medio = receita_total / total
salva = evitaveis * ticket_medio

add_kpi_box(slide, 1.5, 1.8, "Churners evitaveis", f"{evitaveis:,}", color=VERDE)
add_kpi_box(slide, 5.0, 1.8, "Receita preservada", f"R$ {salva:,.0f}", color=VERDE)
add_kpi_box(slide, 8.5, 1.8, "Novo churn", f"{taxa - reducao:.1f}%",
            color=AZUL, delta=f"-{reducao} p.p.", delta_color=VERDE)

add_text(slide, 0.5, 3.8, 12, 1.5,
         f"Uma reducao de {reducao} p.p. no churn preservaria aproximadamente "
         f"R$ {salva:,.0f}\nem receita e {evitaveis:,} clientes.\n\n"
         f"A questao nao e se vale investir em retencao — e onde e como.",
         size=18)

add_text(slide, 0.5, 5.5, 12, 1.5,
         "Os dados apontam para onde olhar:\n"
         f"  - 1o contrato concentra o maior volume de churn ({cr_1o}%)\n"
         f"  - {pct_sil:.0f}% e silencioso — potencialmente recuperavel\n"
         f"  - Dependentes e condicao cronica protegem consistentemente\n"
         f"  - {antecipado_pct:.0f}% dos ativos cancelam com 30+ dias — ha janela\n\n"
         f"A decisao de onde investir e de voces.",
         size=14, color=CINZA)


# ── SLIDE 8: PERGUNTAS ───────────────────────────────────────────
slide = add_slide(prs)
add_title_bar(slide, "Perguntas que os dados levantam")

perguntas = [
    f"O 1o contrato tem {cr_1o}% de churn. O que acontece nos primeiros meses\n"
    f"que leva tantos a nao renovar?",

    f"{pct_sil:.0f}% do churn e silencioso. Quantos desses clientes nem sabiam\n"
    f"que estavam saindo? Sao recuperaveis?",

    f"Dependentes reduzem o churn em ~{cr_sem_dep - cr_com_dep:.0f} p.p. O custo de saida\n"
    f"sobe quando o plano cuida de mais gente — isso e alavancavel?",

    f"Cada cliente que sobrevive ao 1o ciclo vale {clv_2o/clv_1o:.1f}x mais.\n"
    f"Quanto vale investir para garantir essa sobrevivencia?",

    f"R$ {receita_perdida:,.0f} em contratos nao renovados.\n"
    f"Se reduzissemos 5 p.p., preservariamos R$ {salva:,.0f}. Onde comecar?",
]

for i, pergunta in enumerate(perguntas):
    top = 1.5 + i * 1.1
    # Numero
    add_text(slide, 0.5, top, 0.5, 0.5,
             f"{i+1}.", size=20, bold=True, color=AZUL)
    # Texto
    add_text(slide, 1.1, top, 11.5, 0.8, pergunta, size=16, color=AZUL_ESCURO)


# ── SLIDE 9: FECHAMENTO ─────────────────────────────────────────
slide = add_slide(prs)
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = AZUL_ESCURO
shape.line.fill.background()

add_text(slide, 1, 2.5, 11, 1.0,
         "Os dados mostram o caminho.", size=36, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
add_text(slide, 1, 3.8, 11, 0.8,
         "A decisao e de voces.", size=28, color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)
add_text(slide, 1, 5.5, 11, 0.5,
         "Dashboard interativo disponivel para explorar cada segmento em detalhe.",
         size=14, color=CINZA, align=PP_ALIGN.CENTER)


# ── SALVAR ────────────────────────────────────────────────────────
output = "churn_insights.pptx"
prs.save(output)
print(f"\nApresentacao salva em: {output}")
print(f"  9 slides | widescreen 16:9")
print(f"  Dados: {total:,} contratos | {churners_n:,} churners | {taxa}% churn")
